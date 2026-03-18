"""Convert DOCX/XLSX/PPTX/PDF/ODF to text/CSV, then run through MarkdownPreprocessor."""
from __future__ import annotations

import io
import re
from pathlib import Path

from .markdown import MarkdownPreprocessor


def convert(filename: str, content: bytes) -> dict[str, bytes]:
    ext = Path(filename).suffix.lower()
    stem = Path(filename).stem
    fn = {".docx": _docx, ".doc": _docx, ".xlsx": _xlsx, ".xls": _xlsx,
          ".pptx": _pptx, ".ppt": _pptx, ".pdf": _pdf,
          ".odt": _odt, ".ods": _ods, ".odp": _odp}.get(ext)
    return fn(filename, stem, content) if fn else {filename: content}


def _docx(filename: str, stem: str, content: bytes) -> dict[str, bytes]:
    try:
        from docx import Document  # type: ignore
    except ImportError:
        return {filename: f"# {filename}\nInstall python-docx to convert.\n".encode()}
    doc = Document(io.BytesIO(content))
    md, table_no, csvs = [], 0, {}
    for block in doc.element.body:
        tag = block.tag.split("}")[-1]
        if tag == "p":
            for para in doc.paragraphs:
                if para._element is block:
                    text = para.text.strip()
                    style = para.style.name if para.style else ""
                    if not text:
                        md.append("")
                    elif "Heading 1" in style:
                        md.append(f"# {text}")
                    elif "Heading 2" in style:
                        md.append(f"## {text}")
                    elif "Heading 3" in style:
                        md.append(f"### {text}")
                    elif "Heading" in style:
                        m = re.search(r"\d+", style)
                        md.append(f"{'#' * (int(m.group()) if m else 4)} {text}")
                    else:
                        md.append(text)
                    break
        elif tag == "tbl":
            for tbl in doc.tables:
                if tbl._element is block:
                    table_no += 1
                    name = f"{stem}.table_{table_no:03d}.csv"
                    import csv as _csv
                    buf = io.StringIO()
                    w = _csv.writer(buf)
                    for row in tbl.rows:
                        w.writerow([c.text.strip() for c in row.cells])
                    csvs[name] = buf.getvalue().encode()
                    md.append(f"\n[Table {table_no} → `{name}`]\n")
                    break
    # Inventory embedded images (note presence without extracting binary)
    try:
        image_rels = [
            r.target_ref for r in doc.part.rels.values()
            if "image" in r.reltype
        ]
        if image_rels:
            md.append("\n## Embedded Images\n")
            for img_ref in image_rels:
                md.append(f"- `{img_ref}`")
    except Exception:
        pass

    md_name = f"{stem}.md"
    md_bytes = "\n".join(md).encode()
    result = MarkdownPreprocessor().process(md_name, md_bytes)
    result.update(csvs)
    return result


def _xlsx(filename: str, stem: str, content: bytes) -> dict[str, bytes]:
    try:
        import openpyxl  # type: ignore
    except ImportError:
        return {filename: f"# {filename}\nInstall openpyxl to convert.\n".encode()}
    wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
    result, sheet_info = {}, []
    for name in wb.sheetnames:
        ws = wb[name]
        safe = re.sub(r"[^\w\-]", "_", name)
        csv_name = f"{stem}.{safe}.csv"
        import csv as _csv
        buf = io.StringIO()
        w = _csv.writer(buf)
        for row in ws.iter_rows(max_row=ws.max_row, max_col=ws.max_column, values_only=True):
            w.writerow(["" if v is None else str(v) for v in row])
        csv_content = buf.getvalue()
        if csv_content.strip():
            result[csv_name] = csv_content.encode()
            sheet_info.append((name, csv_name))
    idx = [f"# Workbook: {filename}", f"Sheets: {len(sheet_info)}", "",
           "Use `csv_stats` on each sheet. Do NOT use `read_file` on CSVs.", ""]
    for sh, cn in sheet_info:
        idx.append(f"- **{sh}** → `{cn}`")
    result[f"{stem}._index.md"] = "\n".join(idx).encode()
    return result


def _pptx(filename: str, stem: str, content: bytes) -> dict[str, bytes]:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        return {filename: f"# {filename}\nInstall python-pptx to convert.\n".encode()}
    prs = Presentation(io.BytesIO(content))
    md = []
    for no, slide in enumerate(prs.slides, 1):
        md.append(f"## Slide {no}")
        if slide.shapes.title and slide.shapes.title.text:
            md.append(f"**{slide.shapes.title.text.strip()}**")
        md.append("")
        for shape in slide.shapes:
            if not shape.has_text_frame or shape == slide.shapes.title:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    md.append("  " * para.level + ("- " if para.level else "") + text)
        md.append("")
    md_name = f"{stem}.md"
    return MarkdownPreprocessor().process(md_name, "\n".join(md).encode())


def _pdf(filename: str, stem: str, content: bytes) -> dict[str, bytes]:
    try:
        from pdfminer.high_level import extract_pages  # type: ignore
        from pdfminer.layout import LTTextContainer  # type: ignore
    except ImportError:
        return {filename: f"# {filename}\nInstall pdfminer.six to convert.\n".encode()}
    lines = []
    for no, page in enumerate(extract_pages(io.BytesIO(content)), 1):
        lines.append(f"## Page {no}")
        lines.append("")
        for el in page:
            if isinstance(el, LTTextContainer):
                for l in el.get_text().splitlines():
                    if l.rstrip():
                        lines.append(l.rstrip())
        lines.append("")
    txt_name = f"{stem}.txt"
    return MarkdownPreprocessor().process(txt_name, "\n".join(lines).encode())


# ── ODF (OpenDocument) formats ────────────────────────────────────────────────

def _odt(filename: str, stem: str, content: bytes) -> dict[str, bytes]:
    """OpenDocument Text (.odt) → markdown."""
    try:
        from odf import teletype  # type: ignore
        from odf.opendocument import load as odf_load  # type: ignore
        from odf.text import H, List, ListItem, P  # type: ignore
    except ImportError:
        return {filename: f"# {filename}\nInstall odfpy to convert.\n".encode()}

    doc = odf_load(io.BytesIO(content))
    md: list[str] = []

    def _walk(node: object) -> None:
        tag = getattr(node, "qname", ("", ""))[1]
        if tag == "h":
            level = int(node.getAttribute("text:outline-level") or 1)  # type: ignore[attr-defined]
            text = teletype.extractText(node).strip()
            if text:
                md.append(f"{'#' * level} {text}")
        elif tag == "list":
            for item in node.childNodes:  # type: ignore[attr-defined]
                item_text = teletype.extractText(item).strip()
                if item_text:
                    md.append(f"- {item_text}")
        elif tag == "p":
            text = teletype.extractText(node).strip()
            md.append(text)  # empty string → blank line separator
        else:
            for child in getattr(node, "childNodes", []):
                _walk(child)

    for child in doc.text.childNodes:
        _walk(child)

    md_name = f"{stem}.md"
    return MarkdownPreprocessor().process(md_name, "\n".join(md).encode())


def _ods(filename: str, stem: str, content: bytes) -> dict[str, bytes]:
    """OpenDocument Spreadsheet (.ods) → per-sheet CSVs + index."""
    try:
        from odf import teletype  # type: ignore
        from odf.opendocument import load as odf_load  # type: ignore
        from odf.table import Table, TableCell, TableRow  # type: ignore
    except ImportError:
        return {filename: f"# {filename}\nInstall odfpy to convert.\n".encode()}

    import csv as _csv

    doc = odf_load(io.BytesIO(content))
    result: dict[str, bytes] = {}
    sheet_info: list[tuple[str, str]] = []

    for sheet in doc.spreadsheet.getElementsByType(Table):
        name = sheet.getAttribute("table:name") or "Sheet"
        safe = re.sub(r"[^\w\-]", "_", name)
        csv_name = f"{stem}.{safe}.csv"
        buf = io.StringIO()
        w = _csv.writer(buf)
        for row in sheet.getElementsByType(TableRow):
            cells: list[str] = []
            for cell in row.getElementsByType(TableCell):
                repeat = int(cell.getAttribute("table:number-columns-repeated") or 1)
                val = teletype.extractText(cell).strip()
                cells.extend([val] * min(repeat, 256))  # cap repeat to avoid giant sparse sheets
            # Trim trailing empty cells
            while cells and not cells[-1]:
                cells.pop()
            if cells:
                w.writerow(cells)
        csv_content = buf.getvalue()
        if csv_content.strip():
            result[csv_name] = csv_content.encode()
            sheet_info.append((name, csv_name))

    idx = [f"# Spreadsheet: {filename}", f"Sheets: {len(sheet_info)}", "",
           "Use `csv_stats` on each sheet. Do NOT use `read_file` on CSVs.", ""]
    for sh, cn in sheet_info:
        idx.append(f"- **{sh}** → `{cn}`")
    result[f"{stem}._index.md"] = "\n".join(idx).encode()
    return result


def _odp(filename: str, stem: str, content: bytes) -> dict[str, bytes]:
    """OpenDocument Presentation (.odp) → markdown."""
    try:
        from odf import teletype  # type: ignore
        from odf.draw import Page  # type: ignore
        from odf.opendocument import load as odf_load  # type: ignore
        from odf.text import P  # type: ignore
    except ImportError:
        return {filename: f"# {filename}\nInstall odfpy to convert.\n".encode()}

    doc = odf_load(io.BytesIO(content))
    md: list[str] = []

    for no, slide in enumerate(doc.presentation.getElementsByType(Page), 1):
        md.append(f"## Slide {no}")
        for shape in slide.childNodes:
            cls = getattr(shape, "getAttribute", lambda _: None)("presentation:class") or ""
            text = teletype.extractText(shape).strip()
            if not text:
                continue
            if "title" in cls:
                md.append(f"**{text}**")
            else:
                for line in text.splitlines():
                    line = line.strip()
                    if line:
                        md.append(f"- {line}")
        md.append("")

    md_name = f"{stem}.md"
    return MarkdownPreprocessor().process(md_name, "\n".join(md).encode())
